import uuid
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

UPLOADS_DIR = Path(__file__).parent.parent / "uploads"
UPLOADS_DIR.mkdir(exist_ok=True)

# ─── Slide constants (16:9 widescreen) ───────────────────────────────────────
_SW        = Inches(13.33)
_SH        = Inches(7.5)
_HEADER_H  = Inches(0.7)
_HEADER_C  = PRGB(0xDD, 0x0F, 0x2C)
_MG        = Inches(0.4)          # side margin
_CT        = _HEADER_H + Inches(0.3)   # content top
_CH        = _SH - _CT - Inches(0.3)   # content height
_CW        = _SW - _MG * 2              # content width
_WHITE     = PRGB(0xFF, 0xFF, 0xFF)
_DARK      = PRGB(0x1F, 0x23, 0x28)
_GAP       = Inches(0.25)              # gap between columns


def _px(val) -> int:
    return int(val)


# ─── Primitives ───────────────────────────────────────────────────────────────

def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = _SW
    prs.slide_height = _SH
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, fill_color: PRGB, no_border=True):
    sh = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        _px(left), _px(top), _px(width), _px(height),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if no_border:
        sh.line.fill.background()
    return sh


def _add_header(slide, title: str, number: int = 0):
    bar = _rect(slide, 0, 0, _SW, _HEADER_H, _HEADER_C)
    tf = bar.text_frame
    tf.word_wrap = False
    tf.margin_left = _MG
    tf.margin_top = _px(Inches(0))
    tf.margin_right = _MG
    tf.margin_bottom = _px(Inches(0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = (f"{number}. " if number > 0 else "") + title
    run.font.bold = True
    run.font.size = PPt(24)
    run.font.color.rgb = _WHITE


def _bullets(slide, items: list, left, top, width, height):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, dict):
            text = item.get("text", "")
            sub = item.get("sub", "")
        else:
            text = str(item)
            sub = ""
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = PPt(6)
        run = p.add_run()
        run.text = "▸  " + text
        run.font.size = PPt(20)
        run.font.color.rgb = _DARK
        if sub:
            p2 = tf.add_paragraph()
            p2.space_before = PPt(1)
            run2 = p2.add_run()
            run2.text = "      " + sub
            run2.font.size = PPt(15)
            run2.font.color.rgb = PRGB(0x55, 0x55, 0x55)


def _textbox(slide, text: str, left, top, width, height, size=20, bold=False, color=None):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.color.rgb = color or _DARK


def _image(slide, path: str, left, top, width, height):
    p = Path(path)
    if p.exists():
        try:
            from PIL import Image as _PIL
            with _PIL.open(str(p)) as _img:
                iw, ih = _img.size
            if iw > 0 and ih > 0:
                aspect = iw / ih
                tw, th = int(width), int(height)
                if aspect >= tw / th:
                    # ảnh nằm ngang hơn khung — fit theo chiều rộng
                    aw, ah = tw, int(tw / aspect)
                    al, at = int(left), int(top) + (th - int(tw / aspect)) // 2
                else:
                    # ảnh đứng hơn khung — fit theo chiều cao
                    aw, ah = int(th * aspect), th
                    al, at = int(left) + (tw - int(th * aspect)) // 2, int(top)
                slide.shapes.add_picture(str(p), al, at, aw, ah)
                return
        except Exception:
            pass
        slide.shapes.add_picture(str(p), _px(left), _px(top), _px(width), _px(height))
    else:
        ph = _rect(slide, left, top, width, height, PRGB(0xF0, 0xF0, 0xF0), no_border=False)
        ph.line.color.rgb = PRGB(0xCC, 0xCC, 0xCC)
        tf = ph.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run = p2.add_run()
        run.text = "[Hình ảnh]"
        run.font.size = PPt(14)
        run.font.color.rgb = PRGB(0x99, 0x99, 0x99)


_CAP_H = int(Inches(0.28))  # chiều cao dành cho caption bên dưới ảnh


def _image_captioned(slide, path: str, caption: str, left, top, width, height):
    """Render ảnh + caption bên dưới. Nếu caption rỗng thì ảnh chiếm toàn bộ height."""
    if caption:
        img_h = int(height) - _CAP_H
        _image(slide, path, left, top, width, img_h)
        tb = slide.shapes.add_textbox(int(left), int(top) + img_h, int(width), _CAP_H)
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        run.font.size = PPt(11)
        run.font.italic = True
        run.font.color.rgb = PRGB(0x55, 0x55, 0x55)
    else:
        _image(slide, path, left, top, width, height)


def _is_landscape(path: str) -> bool:
    """Trả về True nếu ảnh nằm ngang (width >= height)."""
    try:
        from PIL import Image as _PIL
        with _PIL.open(path) as img:
            w, h = img.size
        return w >= h
    except Exception:
        return True  # mặc định coi là ngang


# ─── Layout templates ─────────────────────────────────────────────────────────

def add_slide_cover(prs: Presentation, title: str, subtitle: str = "") -> None:
    """Cover: full-width red header with large title + optional subtitle below."""
    slide = _blank(prs)
    bar_h = Inches(3.2)
    bar = _rect(slide, 0, 0, _SW, bar_h, _HEADER_C)
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    tf.margin_right = Inches(0.8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = PPt(36)
    run.font.color.rgb = _WHITE
    if subtitle:
        _textbox(slide, subtitle,
                 Inches(0.8), bar_h + Inches(0.4),
                 _SW - Inches(1.6), Inches(1.5),
                 size=20, color=PRGB(0x44, 0x44, 0x44))


def add_slide_bullets(prs: Presentation, title: str, number: int,
                      bullets: list[str]) -> None:
    """Full-width bullet list."""
    slide = _blank(prs)
    _add_header(slide, title, number)
    _bullets(slide, bullets, _MG, _CT, _CW, _CH)


def add_slide_bullets_image1(prs: Presentation, title: str, number: int,
                              bullets: list[str], image_path: str,
                              caption: str = "") -> None:
    """Ảnh ngang: bullets trên (38%) + ảnh dưới (58%). Ảnh dọc: bullets trái (55%) + ảnh phải (40%)."""
    slide = _blank(prs)
    _add_header(slide, title, number)
    if _is_landscape(image_path):
        bh = _px(_CH * 0.38)
        ih = _px(_CH * 0.58)
        gap = _px(_CH) - bh - ih
        _bullets(slide, bullets, _MG, _CT, _CW, bh)
        _image_captioned(slide, image_path, caption, _MG, _px(_CT) + bh + gap, _CW, ih)
    else:
        lw = _px(_CW * 0.55)
        rw = _px(_CW * 0.42)
        _bullets(slide, bullets, _MG, _CT, lw, _CH)
        _image_captioned(slide, image_path, caption, _px(_MG) + lw + _px(_GAP), _CT, rw, _CH)


def add_slide_bullets_image2(prs: Presentation, title: str, number: int,
                              bullets: list[str], image_paths: list[str],
                              captions: list[str] = None) -> None:
    """Ảnh ngang: bullets trên (35%) + 2 ảnh cạnh nhau dưới (60%). Ảnh dọc: bullets trái (50%) + 2 ảnh xếp dọc phải (45%)."""
    slide = _blank(prs)
    _add_header(slide, title, number)
    paths = (image_paths + ["", ""])[:2]
    caps = ((captions or []) + ["", ""])[:2]
    has_landscape = any(_is_landscape(p) for p in paths if p)
    if has_landscape:
        bh = _px(_CH * 0.35)
        ih = _px(_CH * 0.60)
        gap = _px(_CH) - bh - ih
        iw = _px((_CW - _px(_GAP)) / 2)
        _bullets(slide, bullets, _MG, _CT, _CW, bh)
        _image_captioned(slide, paths[0], caps[0], _MG, _px(_CT) + bh + gap, iw, ih)
        _image_captioned(slide, paths[1], caps[1], _px(_MG) + iw + _px(_GAP), _px(_CT) + bh + gap, iw, ih)
    else:
        lw = _px(_CW * 0.50)
        rw = _px(_CW * 0.45)
        img_left = _px(_MG) + lw + _px(_GAP)
        slot_h = _px((_CH - _px(Inches(0.15))) / 2)
        _bullets(slide, bullets, _MG, _CT, lw, _CH)
        _image_captioned(slide, paths[0], caps[0], img_left, _CT, rw, slot_h)
        _image_captioned(slide, paths[1], caps[1], img_left, _px(_CT) + slot_h + _px(Inches(0.15)), rw, slot_h)


def add_slide_images(prs: Presentation, title: str, number: int,
                     image_paths: list[str], captions: list[str] = None) -> None:
    """1 image centered, or 2 images side by side."""
    slide = _blank(prs)
    _add_header(slide, title, number)
    paths = (image_paths + [""])[:2]
    caps = ((captions or []) + ["", ""])[:2]
    if len(image_paths) <= 1:
        iw = _px(_CW * 0.75)
        il = _px(_MG) + _px((_CW - iw) / 2)
        _image_captioned(slide, paths[0], caps[0], il, _CT, iw, _CH)
    else:
        iw = _px((_CW - _px(_GAP)) / 2)
        _image_captioned(slide, paths[0], caps[0], _MG, _CT, iw, _CH)
        _image_captioned(slide, paths[1], caps[1], _px(_MG) + iw + _px(_GAP), _CT, iw, _CH)


def add_slide_image_text(prs: Presentation, title: str, number: int,
                          image_path: str, text: str,
                          caption: str = "") -> None:
    """Image + paragraph text."""
    slide = _blank(prs)
    _add_header(slide, title, number)
    iw = _px(_CW * 0.45)
    tw = _px(_CW - iw - _px(_GAP))
    _image_captioned(slide, image_path, caption, _MG, _CT, iw, _CH)
    _textbox(slide, text, _px(_MG) + iw + _px(_GAP), _CT, tw, _CH, size=16)


# ─── Router ───────────────────────────────────────────────────────────────────

_LAYOUTS = {
    "cover":          lambda prs, d: add_slide_cover(prs, d.get("title", ""), d.get("subtitle", "")),
    "bullets":        lambda prs, d: add_slide_bullets(prs, d["title"], d.get("number", 0), d.get("bullets", [])),
    "bullets_image1": lambda prs, d: add_slide_bullets_image1(prs, d["title"], d.get("number", 0), d.get("bullets", []), (d.get("images") or [""])[0], (d.get("captions") or [""])[0]),
    "bullets_image2": lambda prs, d: add_slide_bullets_image2(prs, d["title"], d.get("number", 0), d.get("bullets", []), d.get("images", []), d.get("captions")),
    "images":         lambda prs, d: add_slide_images(prs, d["title"], d.get("number", 0), d.get("images", []), d.get("captions")),
    "image_text":     lambda prs, d: add_slide_image_text(prs, d["title"], d.get("number", 0), (d.get("images") or [""])[0], d.get("text", ""), (d.get("captions") or [""])[0]),
}


def create_pptx(outline: list[dict], filename: str | None = None) -> str:
    if not filename:
        filename = f"slides_{uuid.uuid4().hex[:8]}.pptx"
    output_path = UPLOADS_DIR / filename
    prs = _new_presentation()
    for slide_data in outline:
        fn = _LAYOUTS.get(slide_data.get("layout", "bullets"), _LAYOUTS["bullets"])
        fn(prs, slide_data)
    prs.save(output_path)
    return str(output_path)


# ─── Word report ──────────────────────────────────────────────────────────────

def write_report_docx(content: str, filename: str | None = None) -> str:
    if not filename:
        filename = f"report_{uuid.uuid4().hex[:8]}.docx"
    output_path = UPLOADS_DIR / filename
    doc = Document()
    doc.styles["Normal"].font.size = Pt(11)
    for line in content.splitlines():
        line = line.rstrip()
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith(("- ", "* ")):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif line.startswith("**") and line.endswith("**"):
            p = doc.add_paragraph()
            p.add_run(line.strip("*")).bold = True
        elif line:
            doc.add_paragraph(line)
    doc.save(output_path)
    return str(output_path)
